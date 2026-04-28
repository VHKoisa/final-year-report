from datetime import date
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE


OUTPUT_FILE = "Final_EndSem_Presentation_Institute_Fixed.pptx"
LOGO_FILE = "Latex Report - Final/svnit.png"
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5

# Theme colors (SVNIT-style conservative academic palette)
NAVY = RGBColor(0x0B, 0x2C, 0x52)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
DARK = RGBColor(0x1F, 0x1F, 0x1F)
MUTED = RGBColor(0x4A, 0x4A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_header_footer(slide):
    # Top branding bar
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_WIDTH_IN), Inches(0.42))
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()

    # Accent strip
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.42), Inches(SLIDE_WIDTH_IN), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()

    # Bottom strip
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.15), Inches(SLIDE_WIDTH_IN), Inches(0.35))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = NAVY
    bottom.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.35), Inches(7.19), Inches(7.6), Inches(0.2))
    p = footer.text_frame.paragraphs[0]
    p.text = "CS402 Industrial Training | Final End-Sem Presentation"
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE

    date_box = slide.shapes.add_textbox(Inches(10.4), Inches(7.19), Inches(2.6), Inches(0.2))
    dp = date_box.text_frame.paragraphs[0]
    dp.text = date.today().strftime("%d %b %Y")
    dp.alignment = PP_ALIGN.RIGHT
    dp.font.name = "Calibri"
    dp.font.size = Pt(11)
    dp.font.color.rgb = WHITE


def add_logo(slide):
    slide.shapes.add_picture(LOGO_FILE, Inches(11.35), Inches(0.55), width=Inches(1.6), height=Inches(0.9))


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.65), Inches(10.3), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = NAVY


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.1), Inches(5.1))
    tf = box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(11)
        p.font.name = "Calibri"
        p.font.size = Pt(24)
        p.font.color.rgb = DARK


def add_two_column(slide, left_title, left_points, right_title, right_points):
    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.7), Inches(4.9))
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
    left.line.color.rgb = NAVY

    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.55), Inches(1.7), Inches(5.7), Inches(4.9))
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(0xFC, 0xF7, 0xF2)
    right.line.color.rgb = ORANGE

    lb = slide.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(5.2), Inches(4.4)).text_frame
    lp = lb.paragraphs[0]
    lp.text = left_title
    lp.font.name = "Calibri"
    lp.font.size = Pt(23)
    lp.font.bold = True
    lp.font.color.rgb = NAVY
    for item in left_points:
        p = lb.add_paragraph()
        p.text = "• " + item
        p.font.name = "Calibri"
        p.font.size = Pt(17)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)

    rb = slide.shapes.add_textbox(Inches(6.75), Inches(1.95), Inches(5.2), Inches(4.4)).text_frame
    rp = rb.paragraphs[0]
    rp.text = right_title
    rp.font.name = "Calibri"
    rp.font.size = Pt(23)
    rp.font.bold = True
    rp.font.color.rgb = ORANGE
    for item in right_points:
        p = rb.add_paragraph()
        p.text = "• " + item
        p.font.name = "Calibri"
        p.font.size = Pt(17)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)


def add_four_quadrants(slide, items):
    labels = ["Why", "Technology", "Business Impact", "Limitations"]
    lefts = [0.8, 6.55, 0.8, 6.55]
    tops = [1.7, 1.7, 4.1, 4.1]
    fills = [RGBColor(0xF5, 0xF8, 0xFC), RGBColor(0xF5, 0xF8, 0xFC), RGBColor(0xFC, 0xF7, 0xF2), RGBColor(0xFC, 0xF7, 0xF2)]

    for i in range(4):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lefts[i]), Inches(tops[i]), Inches(5.7), Inches(2.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fills[i]
        shape.line.color.rgb = NAVY if i < 2 else ORANGE

        tf = slide.shapes.add_textbox(Inches(lefts[i] + 0.2), Inches(tops[i] + 0.12), Inches(5.2), Inches(1.9)).text_frame
        head = tf.paragraphs[0]
        head.text = labels[i]
        head.font.name = "Calibri"
        head.font.size = Pt(19)
        head.font.bold = True
        head.font.color.rgb = NAVY if i < 2 else ORANGE

        for point in items[i]:
            p = tf.add_paragraph()
            p.text = "• " + point
            p.font.name = "Calibri"
            p.font.size = Pt(14)
            p.font.color.rgb = DARK
            p.space_after = Pt(3)


def new_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_footer(slide)
    add_logo(slide)
    add_title(slide, title)
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    # 1. Title
    s = new_slide(prs, "Final End-Sem Presentation")
    sub = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(10.5), Inches(2.4)).text_frame
    lines = [
        "Enterprise Credit Decisioning Modernization, Cost Optimization,",
        "and Fraud Data Centralization",
        "",
        "Prepared by: Vatsal Koisa (U22CS123)",
        "Department of Computer Science and Engineering, SVNIT Surat",
        "Supervised by: Dr. Naveen Kumar",
    ]
    for i, line in enumerate(lines):
        p = sub.paragraphs[0] if i == 0 else sub.add_paragraph()
        p.text = line
        p.font.name = "Calibri"
        p.font.size = Pt(24 if i < 2 else 18)
        p.font.bold = i < 2
        p.font.color.rgb = NAVY if i < 2 else MUTED
        p.space_after = Pt(6)

    # 2. Agenda
    s = new_slide(prs, "Agenda")
    add_bullets(s, [
        "1. Mid-sem recap in one slide",
        "2. Project 1 completion and impact (.adb automation)",
        "3. Project 2 stand-in server: business and technical value",
        "4. Project 3 fraud centralization with Kafka",
        "5. Combined outcomes, limitations, and next steps",
    ])

    # 3. Mid-sem recap slide
    s = new_slide(prs, "Mid-Sem Recap (What Was Already Covered)")
    add_bullets(s, [
        "Identified pain points in fragmented decisioning ecosystem",
        "Explained target architecture for decommissioning and consolidation",
        "Presented manual .adb compile flow and release bottlenecks",
        "Proposed CI-based automated compile and artifact governance plan",
        "Today: continuation from plan to implementation outcomes",
    ])

    # 4. Project 1 continuation
    s = new_slide(prs, "Project 1: Implementation Completed Since Mid-Sem")
    add_bullets(s, [
        "Merged utility repos into central platform to reduce complexity",
        "Enabled automatic .adb compilation on rule changes via CI",
        "Published artifacts to Artifactory with proper versioning",
        "Established traceability and rollback readiness",
        "Reduced manual intervention in business rule releases",
    ])

    # 5. Project 1 analysis
    s = new_slide(prs, "Project 1 Analysis")
    add_four_quadrants(s, [
        [
            "Manual compile flow delayed decisioning rule rollout",
            "Needed reliable, auditable delivery process",
        ],
        [
            "CI-triggered .adb compile pipeline",
            "Artifactory artifact publishing and version control",
        ],
        [
            "Faster policy release and improved compliance posture",
            "Lower human error risk in operational flow",
        ],
        [
            "Pipeline uptime is now business-critical",
            "Legacy cleanup still needs phased execution",
        ],
    ])

    # 6. Project 2 overview
    s = new_slide(prs, "Project 2: In-House Decision Stand-In Server")
    add_two_column(
        s,
        "Why This Was Needed",
        [
            "High recurring cost of cloud and licensed decision engines",
            "Dependency risk when third-party service is delayed/down",
            "Need a controlled internal fallback capability",
        ],
        "Technology and Rollout",
        [
            "Backend: Java Spring Boot",
            "Frontend: React + TypeScript",
            "Backup-first deployment, then parity validation",
            "Gradual traffic shift after business confidence",
        ],
    )

    # 7. Project 2 analysis
    s = new_slide(prs, "Project 2 Analysis")
    add_four_quadrants(s, [
        [
            "Optimize long-term cost and reduce vendor lock-in",
            "Increase continuity for decisioning workloads",
        ],
        [
            "Local deployable internal decision API",
            "Phased adoption strategy with risk control",
        ],
        [
            "Potential multi-million dollar savings",
            "Greater strategic ownership of core capability",
        ],
        [
            "Must match enterprise-grade reliability standards",
            "Requires strict parity testing and governance",
        ],
    ])

    # 8. Project 3 overview
    s = new_slide(prs, "Project 3: Fraud Data Centralization")
    add_two_column(
        s,
        "Problem and Motivation",
        [
            "Identity theft alerts need near real-time action",
            "No central blocked-transaction repository existed",
            "Cross-team inconsistency risk in loan decisions",
        ],
        "Implemented Technical Flow",
        [
            "Producer publishes incoming check requests to Kafka topic",
            "Listeners consume and update central repository",
            "Post-investigation status updates persist centrally",
            "Tech stack: Spring Boot + Kafka",
        ],
    )

    # 9. Project 3 analysis
    s = new_slide(prs, "Project 3 Analysis")
    add_four_quadrants(s, [
        [
            "Fraud handling requires fast and unified visibility",
            "Needed single source of truth for blocked activity",
        ],
        [
            "Event-driven integration using Kafka",
            "Asynchronous producer-consumer architecture",
        ],
        [
            "Faster fraud response and reduced bad approvals",
            "Improved auditability and coordination",
        ],
        [
            "Ordering and deduplication need strong controls",
            "Data privacy and retention governance required",
        ],
    ])

    # 10. Combined impact
    s = new_slide(prs, "Combined End-Sem Outcomes")
    add_bullets(s, [
        "Business: faster policy updates, stronger fraud response, cost optimization path",
        "Technology: automation, event-driven scalability, better system ownership",
        "Governance: improved traceability, reproducibility, and audit readiness",
        "Execution model: practical rollout through phased, low-risk adoption",
    ])

    # 11. Cross-project limitations and controls
    s = new_slide(prs, "Limitations and Risk Controls")
    add_two_column(
        s,
        "Key Limitations",
        [
            "High dependence on pipeline and integration reliability",
            "Parity validation effort for stand-in server migration",
            "Complexity in event ordering and exactly-once handling",
        ],
        "Controls and Mitigations",
        [
            "Monitoring, alerting, and rollback for CI/artifact flow",
            "Pilot-first rollout with measurable acceptance criteria",
            "Idempotency keys, retries, and dead-letter topic handling",
        ],
    )

    # 12. Conclusion
    s = new_slide(prs, "Conclusion and Next Steps")
    add_bullets(s, [
        "Continue quality gates for rule automation pipeline",
        "Run controlled pilots for stand-in server adoption",
        "Production hardening for fraud pipeline observability",
        "Thank you - Q&A",
    ])

    prs.save(OUTPUT_FILE)


if __name__ == "__main__":
    build_deck()
    print(f"Created {OUTPUT_FILE}")
