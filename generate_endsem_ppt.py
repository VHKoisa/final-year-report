from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUTPUT_FILE = "Final_EndSem_Presentation.pptx"


SLIDES = [
    {
        "type": "title",
        "title": "Final End-Sem Presentation",
        "subtitle": "Enterprise Credit Decisioning Modernization, Cost Optimization, and Fraud Data Centralization\nVatsal Koisa | B.Tech CSE | CS402 Industrial Training",
    },
    {
        "type": "bullets",
        "title": "Agenda",
        "bullets": [
            "Mid-sem recap (what was already covered)",
            "Project 1 continuation: Decommissioning + .adb automation",
            "Project 2: In-house stand-in decision server",
            "Project 3: Kafka-based centralized fraud repository",
            "Cross-project impact, limitations, and next steps",
        ],
    },
    {
        "type": "bullets",
        "title": "Mid-Sem Recap (Already Covered)",
        "bullets": [
            "Problem context: fragmented credit decisioning platform",
            "Architecture direction for decommissioning and consolidation",
            "Manual .adb compile flow pain points identified",
            "Proposed plan: automate rule compile and artifact handling",
            "This end-sem focuses on execution outcomes and measurable impact",
        ],
    },
    {
        "type": "bullets",
        "title": "Project 1 Continuation: What Was Implemented",
        "bullets": [
            "Merged utility repositories into a cleaner central codebase",
            "Implemented CI-triggered automatic rule compilation to .adb",
            "Published generated artifacts to Artifactory",
            "Introduced .adb versioning for traceability and rollback",
            "Reduced manual dependency in rule release workflow",
        ],
    },
    {
        "type": "four_block",
        "title": "Project 1 Analysis",
        "why": [
            "Manual rule compile created delays and release risk",
            "Needed better governance for business-critical decisions",
        ],
        "technology": [
            "CI pipeline automation",
            "Artifact repository integration (Artifactory)",
            "Versioned .adb build outputs",
        ],
        "business": [
            "Faster rule rollout to production",
            "Improved auditability and compliance confidence",
            "Lower operational error from manual processing",
        ],
        "limitations": [
            "Pipeline reliability becomes mission-critical",
            "Legacy dependencies still require phased cleanup",
        ],
    },
    {
        "type": "bullets",
        "title": "Project 2: Why Build In-House Stand-In Server",
        "bullets": [
            "Current setup has high recurring cost (cloud + license)",
            "Third-party dependency introduces outage and delay risk",
            "Needed an internal fallback for business continuity",
            "Goal: start as backup, then scale based on parity and trust",
        ],
    },
    {
        "type": "bullets",
        "title": "Project 2: Technology and Rollout",
        "bullets": [
            "Backend: Java Spring Boot",
            "Frontend: React + TypeScript",
            "Architecture objective: local deployable decision stand-in",
            "Rollout strategy: backup-first, parity validation, gradual traffic shift",
        ],
    },
    {
        "type": "four_block",
        "title": "Project 2 Analysis",
        "why": [
            "Reduce vendor lock-in and recurring cost burden",
            "Maintain service continuity during external downtime",
        ],
        "technology": [
            "Internal API-driven decision service",
            "Phased migration with controlled exposure",
        ],
        "business": [
            "Potential multi-million dollar savings",
            "Greater strategic control on a core capability",
        ],
        "limitations": [
            "Must meet enterprise-grade reliability/compliance",
            "Needs robust parity testing vs existing engine",
        ],
    },
    {
        "type": "bullets",
        "title": "Project 3: Fraud Handling Gap",
        "bullets": [
            "Identity theft signals need near-instant propagation",
            "No centralized blocked-transaction repository existed",
            "Cross-team visibility gap could lead to inconsistent decisions",
            "Need a single source of truth for fraud status",
        ],
    },
    {
        "type": "bullets",
        "title": "Project 3: Kafka-Based Centralized Flow",
        "bullets": [
            "Incoming fraud-check requests are produced to Kafka topics",
            "Listeners consume and persist records in central repository",
            "Post-investigation outcomes update the same repository",
            "Stack: Java Spring Boot + Kafka + centralized data store",
        ],
    },
    {
        "type": "four_block",
        "title": "Project 3 Analysis",
        "why": [
            "Fraud response speed directly impacts loss prevention",
            "Teams need synchronized view of blocked/suspected activity",
        ],
        "technology": [
            "Event-driven architecture with Kafka",
            "Asynchronous producer-consumer integration",
        ],
        "business": [
            "Reduced chance of approval during active fraud",
            "Improved traceability and operational coordination",
        ],
        "limitations": [
            "Ordering, deduplication, and retries need strict controls",
            "High governance needs for privacy and retention",
        ],
    },
    {
        "type": "bullets",
        "title": "Combined Outcomes Across All Work",
        "bullets": [
            "Speed: faster change propagation and reduced release cycle time",
            "Reliability: less manual error and more consistent operations",
            "Cost: path toward lower external dependency spend",
            "Control: stronger visibility, governance, and auditability",
        ],
    },
    {
        "type": "bullets",
        "title": "What I Learned",
        "bullets": [
            "Architecture modernization is as much process as code",
            "Event-driven systems require reliability patterns from day one",
            "Business impact framing is essential for technical adoption",
            "Incremental rollout lowers risk in enterprise systems",
        ],
    },
    {
        "type": "bullets",
        "title": "Conclusion and Next Steps",
        "bullets": [
            "Expand quality gates for rule compilation pipeline",
            "Pilot stand-in server in controlled business domains",
            "Production hardening for fraud pipeline observability",
            "Q&A",
        ],
    },
]


def style_title(shape):
    if not shape:
        return
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(34)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x0F, 0x2A, 0x43)


def style_subtitle(shape):
    if not shape:
        return
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def add_bullet_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for idx, bullet in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(0x1F, 0x1F, 0x1F)


def add_four_block_slide(prs, item):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = item["title"]
    style_title(slide.shapes.title)

    labels = ["Why", "Technology", "Business Impact", "Limitations"]
    keys = ["why", "technology", "business", "limitations"]
    lefts = [0.6, 5.0, 0.6, 5.0]
    tops = [1.7, 1.7, 4.0, 4.0]

    for i in range(4):
        box = slide.shapes.add_textbox(Inches(lefts[i]), Inches(tops[i]), Inches(4.1), Inches(2.0))
        tf = box.text_frame
        tf.word_wrap = True

        head = tf.paragraphs[0]
        head.text = labels[i]
        for run in head.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x0F, 0x2A, 0x43)

        for line in item[keys[i]]:
            p = tf.add_paragraph()
            p.text = "• " + line
            for run in p.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(17)
                run.font.color.rgb = RGBColor(0x1F, 0x1F, 0x1F)


def build_presentation():
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = SLIDES[0]["title"]
    title_slide.placeholders[1].text = (
        f"{SLIDES[0]['subtitle']}\nDate: {date.today().strftime('%d %B %Y')}"
    )
    style_title(title_slide.shapes.title)
    style_subtitle(title_slide.placeholders[1])

    for item in SLIDES[1:]:
        if item["type"] == "bullets":
            add_bullet_slide(prs, item["title"], item["bullets"])
        elif item["type"] == "four_block":
            add_four_block_slide(prs, item)

    prs.save(OUTPUT_FILE)


if __name__ == "__main__":
    build_presentation()
    print(f"Created {OUTPUT_FILE}")
